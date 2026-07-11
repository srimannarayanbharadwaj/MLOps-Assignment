"""Generate the video presentation PowerPoint deck with embedded screenshots."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "report" / "VIDEO_PRESENTATION.pptx"

TITLE_COLOR = RGBColor(0, 51, 102)
ACCENT_COLOR = RGBColor(0, 102, 153)
TEXT_COLOR = RGBColor(40, 40, 40)
PLACEHOLDER_FILL = RGBColor(240, 245, 250)
PLACEHOLDER_BORDER = RGBColor(0, 102, 153)

SCREENSHOTS = {
    "eda_class_balance": "screenshots/phase1_eda/class_balance.png",
    "eda_histograms": "screenshots/phase1_eda/numeric_feature_histograms.png",
    "eda_heatmap": "screenshots/phase1_eda/correlation_heatmap.png",
    "mlflow_experiments": "screenshots/phase3_mlflow/01_experiments_page.png",
    "mlflow_runs": "screenshots/phase3_mlflow/02_run_training_table.png",
    "mlflow_logreg": "screenshots/phase3_mlflow/02_run_logreg_table.png",
    "mlflow_artifacts": "screenshots/phase3_mlflow/04_logistic_artifacts.png",
    "github_actions": "screenshots/phase3_mlflow/Screenshot 2026-07-11 214818.png",
    "docker_run": "screenshots/phase6_docker/docker_run.png",
    "docker_predict": "screenshots/phase6_docker/03_curl_predict_success.png",
    "k8s_minikube": "screenshots/phase7_k8s/minicube start success.png",
    "k8s_deploy": "screenshots/phase7_k8s/k8s deployment.png",
    "k8s_port_forward": "screenshots/phase7_k8s/k8s port-forward.png",
    "k8s_health": "screenshots/phase7_k8s/k8s health success.png",
    "monitoring_logs": "screenshots/phase8_monitoring/Request logs.png",
    "monitoring_metrics": "screenshots/phase8_monitoring/metrics endpoint.png",
}


def set_run_font(run, size: int = 18, bold: bool = False, color=TEXT_COLOR) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def resolve_images(keys: list[str]) -> list[Path]:
    paths = []
    for key in keys:
        rel = SCREENSHOTS.get(key)
        if not rel:
            continue
        full = ROOT / rel
        if full.exists():
            paths.append(full)
    return paths


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Heart Disease MLOps"
    set_run_font(run, size=40, bold=True, color=TITLE_COLOR)
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(8.8), Inches(2.5))
    stf = sub.text_frame
    lines = [
        "End-to-End Machine Learning Operations Pipeline",
        "",
        "Sriman Narayan Bharadwaj R",
        "MLOps Assignment 01 (2026)",
        "github.com/srimannarayanbharadwaj/MLOps-Assignment",
    ]
    for idx, line in enumerate(lines):
        para = stf.paragraphs[0] if idx == 0 else stf.add_paragraph()
        run = para.add_run()
        run.text = line
        size = 22 if idx == 0 else 18
        set_run_font(run, size=size, bold=idx == 0, color=ACCENT_COLOR if idx == 0 else TEXT_COLOR)
        para.alignment = PP_ALIGN.CENTER

    slide.notes_slide.notes_text_frame.text = (
        "Introduce yourself, assignment name, and repository link."
    )


def add_section_header(slide, title: str, subtitle: str = "") -> None:
    header = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.7))
    hp = header.text_frame.paragraphs[0]
    run = hp.add_run()
    run.text = title
    set_run_font(run, size=28, bold=True, color=TITLE_COLOR)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.4), Inches(0.85), Inches(9.2), Inches(0.4))
        sp = sub.text_frame.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        set_run_font(srun, size=14, color=ACCENT_COLOR)


def add_bullets(slide, bullets: list[str], left=0.4, top=1.35, width=4.5, height=4.8) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = bullet
        para.level = 0
        para.space_after = Pt(8)
        for run in para.runs:
            set_run_font(run, size=16)


def add_screenshot_placeholder(
    slide,
    label: str,
    left: float = 5.1,
    top: float = 1.35,
    width: float = 4.4,
    height: float = 4.8,
) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PLACEHOLDER_FILL
    shape.line.color.rgb = PLACEHOLDER_BORDER
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Missing Screenshot\n\n" + label
    set_run_font(run, size=14, bold=True, color=ACCENT_COLOR)


def add_images_panel(
    slide,
    image_keys: list[str],
    left: float = 5.0,
    top: float = 1.35,
    width: float = 4.5,
    height: float = 5.0,
    fallback_label: str = "Screenshot",
) -> None:
    images = resolve_images(image_keys)
    if not images:
        add_screenshot_placeholder(slide, fallback_label, left, top, width, height)
        return

    gap = 0.08
    slot_height = (height - gap * (len(images) - 1)) / len(images)
    for index, image_path in enumerate(images):
        y_pos = top + index * (slot_height + gap)
        slide.shapes.add_picture(
            str(image_path),
            Inches(left),
            Inches(y_pos),
            width=Inches(width),
            height=Inches(slot_height),
        )


def add_content_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    bullets: list[str],
    image_keys: list[str],
    speaker_notes: str,
    fallback_label: str = "Screenshot",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, title, subtitle)
    add_bullets(slide, bullets)
    add_images_panel(slide, image_keys, fallback_label=fallback_label)
    slide.notes_slide.notes_text_frame.text = speaker_notes


def add_full_width_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    speaker_notes: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, title)
    add_bullets(slide, bullets, left=0.6, top=1.3, width=8.8, height=5.5)
    slide.notes_slide.notes_text_frame.text = speaker_notes


def add_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "System Architecture", "End-to-end MLOps flow")

    flow = (
        "UCI CSV\n"
        "   |\n"
        "data_download.py -> cleaned dataset\n"
        "   |\n"
        "eda.py + train.py -> MLflow + metrics\n"
        "   |\n"
        "package_model.py -> joblib pipeline\n"
        "   |\n"
        "FastAPI -> Docker -> Kubernetes\n"
        "   |\n"
        "/predict  /health  /metrics\n\n"
        "GitHub Actions + pytest validate each stage"
    )
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(4.5), Inches(5.0))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = flow
    set_run_font(run, size=15)
    run.font.name = "Consolas"

    add_images_panel(
        slide,
        ["mlflow_experiments", "k8s_deploy"],
        left=5.2,
        top=1.4,
        width=4.3,
        height=5.0,
        fallback_label="Architecture overview",
    )
    slide.notes_slide.notes_text_frame.text = (
        "Explain how each component connects from raw data to deployed API."
    )


def add_results_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Model Comparison Results", "5-fold Stratified Cross-Validation")

    rows, cols = 3, 5
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.4), Inches(4.6), Inches(1.3))
    table = table_shape.table

    headers = ["Model", "Accuracy", "Precision", "Recall", "ROC-AUC"]
    data = [
        ["Logistic Regression", "0.8447", "0.8475", "0.8050", "0.9188"],
        ["Random Forest", "0.8416", "0.8388", "0.8127", "0.9163"],
    ]

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                set_run_font(run, size=12, bold=True, color=RGBColor(255, 255, 255))
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    set_run_font(run, size=12, bold=col_idx == 0)

    add_bullets(
        slide,
        [
            "Selected model: Logistic Regression",
            "Reason: highest ROC-AUC (0.9188)",
            "Random Forest had slightly higher recall",
            "Logistic Regression is easier to explain and deploy",
        ],
        left=0.5,
        top=3.0,
        width=4.6,
        height=2.2,
    )

    add_images_panel(
        slide,
        ["mlflow_logreg", "mlflow_artifacts"],
        left=5.2,
        top=1.4,
        width=4.3,
        height=5.0,
        fallback_label="MLflow run metrics and artifacts",
    )
    slide.notes_slide.notes_text_frame.text = (
        "Walk through the metrics table and justify model selection."
    )


def add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Conclusion", "Project delivered end-to-end")

    bullets = [
        "Built full MLOps pipeline on UCI Cleveland Heart Disease data",
        "Tracked experiments with MLflow and packaged the best model",
        "Validated with pytest and GitHub Actions CI/CD",
        "Deployed via Docker, Minikube/Kubernetes, and monitoring",
        "Repository: github.com/srimannarayanbharadwaj/MLOps-Assignment",
        "Thank you for watching",
    ]
    add_bullets(slide, bullets, left=0.8, top=1.5, width=8.4, height=4.5)
    slide.notes_slide.notes_text_frame.text = (
        "Summarize all phases and repeat the GitHub repository link."
    )


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_full_width_content_slide(
        prs,
        "Agenda",
        [
            "1. Project setup and repository structure",
            "2. Data download and EDA",
            "3. Model training and selection",
            "4. MLflow experiment tracking",
            "5. Model packaging",
            "6. Unit tests and CI/CD",
            "7. Docker containerization",
            "8. Kubernetes deployment (Minikube)",
            "9. Monitoring and conclusion",
        ],
        "Give a quick roadmap of what you will demonstrate live.",
    )

    add_content_slide(
        prs,
        "Phase 0: Project Setup",
        "Environment and repository structure",
        [
            "Python 3.12 virtual environment (.venv312)",
            "Pinned dependencies in requirements.txt",
            "Folders: data, src, tests, models, k8s, screenshots",
            "GitHub repo stores code, tests, CI, and deployment files",
            "Goal: reproducible MLOps workflow from clean clone",
        ],
        ["mlflow_experiments"],
        "Show project folders live and mention Python 3.12 requirement.",
        "MLflow experiments overview",
    )

    add_content_slide(
        prs,
        "Phase 1: Data and EDA",
        "UCI Cleveland Heart Disease dataset",
        [
            "303 rows, 13 clinical features + target",
            "Missing '?' values handled in ca and thal",
            "Target binarized: 0 = no disease, 1+ -> 1",
            "Class balance: 164 vs 139",
            "EDA: histograms, correlation heatmap, class chart",
        ],
        ["eda_class_balance", "eda_histograms", "eda_heatmap"],
        "Run data_download and explain what each EDA plot shows.",
        "EDA plots",
    )

    add_content_slide(
        prs,
        "Phase 2: Feature Engineering and Models",
        "Preprocessing + model comparison",
        [
            "Numeric features -> StandardScaler",
            "Categorical features -> OneHotEncoder",
            "Models: Logistic Regression vs Random Forest",
            "Evaluation: 5-fold stratified cross-validation",
            "Metrics: accuracy, precision, recall, ROC-AUC",
        ],
        ["mlflow_runs"],
        "Show preprocessing pipeline and cross-validation results.",
        "MLflow run comparison",
    )

    add_results_slide(prs)

    add_content_slide(
        prs,
        "Phase 3: MLflow Tracking",
        "Experiment logging and comparison",
        [
            "Logged params, metrics, and artifacts per run",
            "Compared logistic_regression vs random_forest",
            "Saved confusion matrix and sklearn pipeline",
            "MLflow UI used for experiment comparison",
            "Experiment: heart-disease-classification",
        ],
        ["mlflow_experiments", "mlflow_runs"],
        "Start mlflow ui and show runs table and artifacts.",
        "MLflow experiments and runs",
    )

    add_content_slide(
        prs,
        "Phase 4: Packaging",
        "Reproducible model artifact",
        [
            "Saved full pipeline with joblib",
            "Artifact: heart_disease_pipeline.joblib",
            "Metadata: model_metadata.json",
            "Sample input: sample_input.json",
            "Fresh clone + requirements.txt reproduces model",
        ],
        ["mlflow_artifacts"],
        "Run package_model and show saved artifacts.",
        "Packaged model artifacts",
    )

    add_content_slide(
        prs,
        "Phase 5: Tests and CI/CD",
        "Automated quality checks on every push",
        [
            "test_data.py: null handling and target encoding",
            "test_model.py: prediction shape and probabilities",
            "test_api.py: /health and /metrics endpoints",
            "GitHub Actions: lint -> pytest -> training job",
            "CI ensures reproducibility and catches regressions",
        ],
        ["github_actions"],
        "Run pytest locally and show green Actions run in browser.",
        "GitHub Actions workflow",
    )

    add_content_slide(
        prs,
        "Phase 6: Docker + FastAPI",
        "Containerized model serving",
        [
            "FastAPI endpoints: /health, /predict, /metrics",
            "Dockerfile builds Python 3.12 image with model",
            "docker build and docker run expose port 8000",
            "Live prediction returns class + confidence",
            "Example: prediction=0, confidence=0.8292",
        ],
        ["docker_run", "docker_predict"],
        "Build image, run container, and call /predict live.",
        "Docker run and predict response",
    )

    add_content_slide(
        prs,
        "Phase 7: Kubernetes (Minikube)",
        "Local cluster deployment",
        [
            "Deployment + Service YAML in k8s/",
            "minikube image load imports local Docker image",
            "kubectl apply deploys API to cluster",
            "Readiness/liveness probes use /health",
            "port-forward exposes service locally",
        ],
        ["k8s_minikube", "k8s_deploy", "k8s_health"],
        "Show pod Running 1/1 and successful health response.",
        "Minikube and Kubernetes deployment",
    )

    add_content_slide(
        prs,
        "Phase 8: Monitoring",
        "Logging and Prometheus metrics",
        [
            "Middleware logs timestamp, endpoint, latency",
            "Logs help debug slow or failing requests",
            "/metrics exposes Prometheus text format",
            "Tracks request count and duration histograms",
            "Supports production observability basics",
        ],
        ["monitoring_logs", "monitoring_metrics"],
        "Start uvicorn, send requests, show logs and metrics.",
        "Request logs and metrics endpoint",
    )

    add_architecture_slide(prs)

    add_full_width_content_slide(
        prs,
        "Challenges and Lessons Learned",
        [
            "Python 3.14 broke MLflow UI; Python 3.12 fixed it",
            "Docker and Minikube port conflicts required alternate ports",
            "Pinned requirements improved reproducibility across machines",
            "MLflow simplified experiment comparison and artifact tracking",
            "CI pipeline gave confidence that changes did not break the workflow",
        ],
        "Mention real issues you faced and how you solved them.",
    )

    add_closing_slide(prs)

    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")
    for key, rel in SCREENSHOTS.items():
        status = "ok" if (ROOT / rel).exists() else "missing"
        print(f"  [{status}] {rel}")


if __name__ == "__main__":
    main()
